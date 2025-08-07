from pptx import Presentation
import streamlit as st
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain.prompts import PromptTemplate
from langchain.chains.summarize import load_summarize_chain
from langchain.retrievers import ContextualCompressionRetriever
from langchain.retrievers.document_compressors import LLMChainExtractor
# from langchain.chains import LLMChain
from langchain_community.vectorstores import FAISS
# from langchain_huggingface import HuggingFaceEmbeddings
from langchain_openai import OpenAIEmbeddings
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain

def load_and_chunk_pdf(pdf_file, chunk_size=1000, chunk_overlap=200,temp_filename: str = "temp.pdf"):
    """
    Carrega um arquivo PDF e o converte em uma lista de documentos segmentados, com metadados por página.

    Este método salva temporariamente o PDF, extrai o conteúdo de cada página, concatena o texto completo para
    ter continuidade de idéias através das páginas e o divide em trechos (chunks) de tamanho definido. Cada trecho 
    é enriquecido com informações da página original e convertido em um objeto `Document` da LangChain.
    """
    pdf_file.seek(0)
    # Save uploaded PDF to a temporary file
    with open(temp_filename, "wb") as f:
        f.write(pdf_file.read())

    # Load pages
    loader = PyPDFLoader(temp_filename)
    pages = loader.load()
    
    # Prepare to track page indices
    full_text = ""
    page_start_indices = []
    current_index = 0

    for page_num, page in enumerate(pages, start=1):
        page_start_indices.append({
            "page_number": page_num,
            "start_index": current_index
        })
        full_text += page.page_content + "\n"
        current_index += len(page.page_content) + 1  # +1 for newline

    # Split text into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    chunks = splitter.split_text(full_text)

    # Attach metadata
    file_name = pdf_file.name
    documents = []

    for chunk_text in chunks:
        chunk_start_index = full_text.find(chunk_text)
        start_page_number = 1
        for page_info in reversed(page_start_indices):
            if chunk_start_index >= page_info["start_index"]:
                start_page_number = page_info["page_number"]
                break

        text_with_page = f"(Page {start_page_number})\n{chunk_text}"
        doc = Document(
            page_content=text_with_page,
            metadata={
                "page_number": start_page_number,
                "file_name": file_name
            }
        )
        documents.append(doc)

    return documents, len(pages)

def generate_summary_from_documents(documents, llm):
    """
    Gera um resumo clínico estruturado a partir de uma lista de documentos segmentados, utilizando um modelo de linguagem.

    A função aplica uma abordagem de dois estágios ("map-reduce"): primeiro resume individualmente cada trecho do texto,
    focando em informações clínicas relevantes, e depois combina esses resumos em uma explicação final clara, técnica
    e em português brasileiro.
    """
    # Prompt for individual chunk summarization
    map_prompt = PromptTemplate(
        input_variables=["text"],
        template="""
You are a medical research assistant trained to extract the most clinically important details from medical literature.

Your task is to summarize the following text chunk by focusing strictly on the **clinically relevant information** only.

🎯 What to include:
- Key clinical findings
- Diagnostic strategies or criteria
- Treatment or management approaches
- Study context or outcomes relevant to patient care
- Pathophysiological insights if they influence clinical decision-making

⚠️ Strict Instructions:
- Do NOT include titles, headings, or formatting (e.g. bold or bullet points).
- Do NOT speculate or generalize — stay factual.
- Focus on clinical findings, patient management, study context, or implications.
- If the chunk refers to tables, mention only their overall clinical meaning (no row-level data).

Text:
{text}

Summary:
"""
    )

    # Prompt for final combined summary in Portuguese
    reduce_prompt = PromptTemplate(
        input_variables=["text"],
        template="""
You are a medical professor giving a private one-on-one lesson to a medical student preparing for an oral exam. Your job is not only to summarize the article, but to explain it in a clear, structured, and clinically useful way — using only what is present in the text.

The explanation must be written in **Brazilian Portuguese**, using accessible but technically correct language. Think of a professor guiding a student preparing for high-stakes clinical reasoning.

Use the following structure in your response:

---

1. **Introdução (Introduction):**
- Briefly describe the main clinical problem the article addresses.
- Why is this topic relevant to medical practice?
- What should the student aim to learn and retain?

2. **Métodos (Study Methods):**
- If the article presents original research, describe the research design (e.g., cohort study, retrospective analysis, review).
- If it's a guideline or clinical review, describe the type of approach used (e.g., literature review, case discussions).
- Do **not confuse diagnostic tests with research methodology.**

3. **Resultados (Findings):**
Carefully explain the most important findings and data. When present in the text, include:
- **Quantitative cutoffs or thresholds** (e.g., lab reference values, decision criteria, score ranges)
- **Named diagnostic tests or procedures** (e.g., SPEP, FLC, immunofixação, CT scan)
- **Decision-making logic** (e.g., "when to biopsy", "what defines low-risk")
- **Risk models** or stratification frameworks (e.g., clinical scores, staging)
- **Follow-up schedules** (e.g., "monitor every 6 months with X and Y")
- **Associations or complications** (e.g., neuropathy, renal disease, risk of transformation)
- **Genetic or familial predisposition**, when mentioned (e.g., parentes de 1º grau com MGUS ou mieloma)
- Use phrases like “isso significa que...” to explain clinical relevance.

4. **Conclusão (Conclusion):**
- Summarize the core clinical message.
- What change in reasoning or decision-making should come from this article?

---

📌 If relevant, also include additional sections such as:
- *"Implicações clínicas"*
- *"Discussão de casos"*
- *"Desafios diagnósticos"*
- *"Aspectos psicológicos"*
- *"Impacto econômico"*

🎯 Instructions:
- When present in the text, always include relevant quantitative data (e.g., cutoffs, thresholds, prevalence rates) to substantiate your summary and enhance clinical precision.
- Do **not invent** or infer beyond the provided text.
- Do **not copy** sentences directly — rephrase and teach.
- Emphasize **reasoning, clarity, and retention**.
- Prioritize information that is clinically actionable.

Here are the chunk summaries:
{text}
"""
    )

    # Run summarization pipeline
    chain = load_summarize_chain(
        llm,
        chain_type="map_reduce",
        map_prompt=map_prompt,
        combine_prompt=reduce_prompt
    )

    result = chain.invoke(documents)
    return result["output_text"]

def generate_slide_markdown_from_summary(summary_text: str, llm) -> str:
    """
    Converte um resumo clínico em formato narrativo para uma estrutura de apresentação em Markdown, organizada por slides com tópicos.

    Utiliza um modelo de linguagem para transformar o conteúdo em blocos temáticos claros e didáticos,
    com foco na comunicação eficiente de achados clínicos para uso em apresentações.
    """
    # Prompt to convert final_summary into bullet-style slides (in Portuguese)
    slide_bullet_prompt = PromptTemplate(
    input_variables=["summary"],
    template="""
Você é um assistente especializado em criar apresentações de PowerPoint para profissionais de saúde.

Sua tarefa é transformar o resumo abaixo em uma apresentação clara, concisa e didática, estruturada em formato de slides com tópicos.


🎯 Objetivo:
Gerar uma apresentação que resuma os principais pontos clínicos do conteúdo, permitindo que estudantes ou profissionais compreendam rapidamente os achados, condutas e implicações do artigo.

🧱 Instruções:
- Divida o conteúdo em múltiplos slides.
- Cada slide deve conter:
  - Um título representando o bloco temático principal (por exemplo: "Definição", "Critérios Diagnósticos", "Estratégias de Manejo", "Complicações", "Implicações Clínicas", etc.).
  - Entre 3 e 5 tópicos por slide, com no máximo 2 linhas cada.
- Use linguagem clara, objetiva e didática.
- Se uma seção tiver muito conteúdo, divida em mais de um slide com o mesmo título (sem numerar, nem usar “1/2”, “2/2”…).
- Tente sempre utilizar todos os dados quantitativos, pois são importantes para apresentação
- Não invente nenhuma informação.
- Use o formato Markdown:
  - Comece cada slide com `### Título do Slide`
  - Escreva os tópicos com `- ` no início de cada linha.
- **Não copie parágrafos inteiros** nem repita frases longas.
- **Nâo utilize a linguagem de aluno e professor se presente no resumo para criar os conteúdos, se atente aos fatos

Resumo:
{summary}
"""
)
    slide_generator = slide_bullet_prompt | llm
    return slide_generator.invoke({"summary": summary_text}).content


def export_markdown_to_pptx(markdown_text: str, output_path: str = "presentation.pptx"):
    """
    Converte um texto em Markdown estruturado por slides em um arquivo de apresentação PowerPoint (.pptx).

    Cada seção iniciada por `###` é transformada em um slide com título e tópicos, organizados de forma clara e formatados para fins didáticos.
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]

    sections = markdown_text.split("### ")
    for section in sections:
        if not section.strip():
            continue
        try:
            title, content = section.strip().split("\n", 1)
        except ValueError:
            continue

        bullets = [line.strip("- ").strip() for line in content.strip().split("\n") if line.strip()]
        chunks = [bullets[i:i + 5] for i in range(0, len(bullets), 5)]

        for chunk in chunks:
            slide = prs.slides.add_slide(layout)
            title_shape = slide.shapes.title
            title_shape.text_frame.clear()
            p = title_shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = title.strip()
            run.font.name = 'Times New Roman'
            run.font.size = Pt(40)

            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            for bullet in chunk:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.name = 'Times New Roman'
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.JUSTIFY
                p.level = 0

    prs.save(output_path)

def build_compressed_retriever(documents, llm):
    """
    Cria um retriever comprimido a partir de documentos clínicos, usando embeddings e um compressor baseado em LLM.

    O objetivo é otimizar a recuperação de informações clínicas relevantes, reduzindo o ruído e mantendo somente o conteúdo essencial para responder perguntas médicas com base no texto original.
    """
    # Step 1: Embedding model (Hugging Face)
    # embedding_model = HuggingFaceEmbeddings(
    #     model_name="sentence-transformers/all-MiniLM-L6-v2"
    # )
    embedding_model = OpenAIEmbeddings(
        model="text-embedding-3-small",
        openai_api_key=llm.openai_api_key
    )

    # Step 2: Create FAISS vectorstore
    vectorstore = FAISS.from_documents(documents, embedding_model)

    # Step 3: Base retriever
    base_retriever = vectorstore.as_retriever(
        search_type="similarity",
        search_kwargs={"k": 5}
    )

    # Step 4: Compressor prompt
    compressor_prompt = PromptTemplate(
        input_variables=["context"],
        template="""
You are a medical information extractor.

Your task is to compress and extract only the clinically **relevant and numerically precise** information from the provided document.

✅ Always preserve:
- **Page number** (e.g. “(Page 3)”) exactly as it appears.
- **Numerical values** (percentages, cutoff thresholds, rates, time intervals, test results).
- **Diagnostic and prognostic criteria**.

💡 Include:
- All risk scores, percentages, stratification models (e.g. "progression risk is 1% per year").
- Differences between disease subtypes.
- Any statistical comparison or quantitative data.

🚫 Ignore:
- Bibliography or article metadata.
- Mentions of references unless clinically relevant.

Here is the text for compression:

{context}

"""
    )

    # Step 5: Build compressor
    compressor = LLMChainExtractor.from_llm(llm=llm, prompt=compressor_prompt)

    # Step 6: Wrap with compression retriever
    retriever = ContextualCompressionRetriever(
        base_retriever=base_retriever,
        base_compressor=compressor
    )

    return retriever    

def build_qa_chain(llm, retriever):
    """
    Constrói uma cadeia de perguntas e respostas conversacional, baseada em documentos clínicos e com memória de contexto.

    A cadeia permite interações contínuas com o modelo, mantendo o histórico da conversa e gerando respostas em português,
    sempre com base no conteúdo do artigo e indicando os números das páginas como referência.
    """ 
    prompt_template = """
You are a medical research assistant.

Here is the previous conversation for context:
{chat_history}

You are having an ongoing conversation with a medical professional. Use the previous conversation naturally to understand and respond to follow-up questions, as if you are talking to a colleague, ensuring clarity and continuity.

Using only the provided document context, answer the following question clearly and concisely for a Brazilian medical professional in Portuguese, even if the document context is in English.
If there is any data to confirm your analysis in the context, try to use it.

For each paragraph or major piece of information in your answer, indicate the corresponding page number at the end of the paragraph in parentheses. For example:

"Texto da resposta. (página 5)"

When providing references, always cite the page number(s) from the document context instead of article citation numbers. If the retrieved context does not explicitly include a page number, write (página não especificada).

If the context does not contain enough information to answer, don't invent anything.

Document Context:
{context}

Question:
{question}

Answer:
"""

    prompt = PromptTemplate(
        input_variables=["chat_history", "context", "question"],
        template=prompt_template
    )

    memory = ConversationBufferMemory(
        memory_key="chat_history",
        return_messages=True
    )

    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        memory=memory,
        combine_docs_chain_kwargs={"prompt": prompt}
    )

    return qa_chain

def translate_to_english(question_pt, llm):
    prompt_translation = f"""
You are a professional translator. Translate the following text to English, preserving its medical meaning precisely,
also only translate, don't add any new information or idea.

{question_pt}
"""
    return llm.invoke(prompt_translation).content